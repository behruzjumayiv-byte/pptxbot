import logging
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import textwrap

logger = logging.getLogger(__name__)

class PPTMaker:
    
    
    def __init__(self):
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
    
    async def create_presentation(self, topic: str, author: str, slides_content: list, template_num: int) -> str:
        """
        PPTX prezentatsiya yaratish
        
        Args:
            topic: Mavzu
            author: Muallif
            slides_content: Slaydlar ma'lumotlari
            template_num: Shablon raqami
            
        Returns:
            str: Yaratilgan fayl yo'li
        """
        try:
            # Prezentatsiya yaratish
            prs = Presentation()
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height
            
            # Shablon yo'li
            template_path = f"templates/{template_num}.png"
            
            if not os.path.exists(template_path):
                logger.warning(f"Shablon topilmadi: {template_path}")
                template_path = None
            
            # Har bir slayd uchun
            for idx, slide_data in enumerate(slides_content):
                title = slide_data.get('title', 'Untitled')
                content = slide_data.get('content', '')
                
                # Bo'sh slayd yaratish
                slide_layout = prs.slide_layouts[6]  # Bo'sh layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Shablon fonini qo'shish
                if template_path and os.path.exists(template_path):
                    try:
                        left = Inches(0)
                        top = Inches(0)
                        pic = slide.shapes.add_picture(
                            template_path,
                            left, top,
                            width=self.slide_width,
                            height=self.slide_height
                        )
                        # Fonni orqaga surish
                        slide.shapes._spTree.remove(pic._element)
                        slide.shapes._spTree.insert(2, pic._element)
                    except Exception as e:
                        logger.error(f"Shablon qo'shishda xatolik: {e}")
                
                # Sarlavha qo'shish
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5),
                    Inches(9), Inches(1.2)
                )
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
                
                # Kontent qo'shish
                if content:
                    # Matnni avtomatik wrap qilish
                    wrapped_content = self._wrap_text(content, 80)
                    
                    content_box = slide.shapes.add_textbox(
                        Inches(0.8), Inches(2),
                        Inches(8.4), Inches(4.5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True
                    
                    p = content_frame.paragraphs[0]
                    p.text = wrapped_content
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.LEFT
                    p.line_spacing = 1.3
                    
                    # Agar matn juda uzun bo'lsa, shrift o'lchamini kamaytirish
                    if len(wrapped_content) > 500:
                        p.font.size = Pt(14)
                    elif len(wrapped_content) > 300:
                        p.font.size = Pt(16)
            
            # Faylni saqlash
            output_filename = f"{self._sanitize_filename(topic)}.pptx"
            output_path = os.path.join("output", output_filename)
            
            # Output papkasini yaratish
            os.makedirs("output", exist_ok=True)
            
            prs.save(output_path)
            logger.info(f"Prezentatsiya saqlandi: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"PPTX yaratishda xatolik: {e}")
            raise
    
    def _wrap_text(self, text: str, width: int) -> str:
        """Matnni avtomatik wrap qilish"""
        try:
            lines = text.split('\n')
            wrapped_lines = []
            
            for line in lines:
                if len(line) <= width:
                    wrapped_lines.append(line)
                else:
                    wrapped = textwrap.wrap(line, width=width)
                    wrapped_lines.extend(wrapped)
            
            return '\n'.join(wrapped_lines)
        except Exception as e:
            logger.error(f"Matn wrap xatolik: {e}")
            return text
    
    def _sanitize_filename(self, filename: str) -> str:
        
        # Ruxsat etilmagan belgilarni olib tashlash
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '')
        
        # Bo'sh joylarni chiziqcha bilan almashtirish
        filename = filename.replace(' ', '_')
        
        # Maksimal uzunlik
        if len(filename) > 50:
            filename = filename[:50]
        
        return filename
